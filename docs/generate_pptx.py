"""
Genera DocuMind_Pitch_Deck.pptx — 15 slides, diseño oscuro minimalista futurista.

Uso:
  cd "DocuMind AI"
  python docs/generate_pptx.py

Requiere: pip install python-pptx matplotlib numpy
"""

import os
import io

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

OUTPUT = os.path.join(os.path.dirname(__file__), "DocuMind_Pitch_Deck.pptx")

# ── Paleta ────────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD   = RGBColor(0x1E, 0x29, 0x3B)
BG_CARD2  = RGBColor(0x27, 0x34, 0x4D)
VIOLET    = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_LT = RGBColor(0xA7, 0x8B, 0xFA)
EMERALD   = RGBColor(0x10, 0xB9, 0x81)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LT   = RGBColor(0x94, 0xA3, 0xB8)
GRAY_DK   = RGBColor(0x47, 0x55, 0x69)
RED_SOFT  = RGBColor(0xEF, 0x44, 0x44)

# Hex para matplotlib
_BG   = "#0F172A"
_CARD = "#1E293B"
_VIO  = "#7C3AED"
_VIOL = "#A78BFA"
_EME  = "#10B981"
_AMB  = "#F59E0B"
_GRY  = "#94A3B8"
_RED  = "#EF4444"
_WHT  = "#FFFFFF"

W = Inches(13.33)
H = Inches(7.5)
M = Inches(0.5)

RND  = 5   # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
FLAT = 1   # MSO_AUTO_SHAPE_TYPE.RECTANGLE


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _new_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return s


def _rect(slide, x, y, w, h, fill=BG_CARD, line_color=None,
          rounded=True, radius=0.07):
    shape = slide.shapes.add_shape(RND if rounded else FLAT, x, y, w, h)
    if rounded:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _txb(slide, x, y, w, h, text="", size=18, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.text_frame.word_wrap = True
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _add_para(txb, text, size=16, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = txb.text_frame.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def _badge(slide, x, y, w, h, text, bg=VIOLET, tc=WHITE, size=10):
    shape = slide.shapes.add_shape(RND, x, y, w, h)
    shape.adjustments[0] = 0.45          # pill shape
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = tc


def _slide_title(slide, title, subtitle=None, color=VIOLET_LT):
    _txb(slide, M, Inches(0.28), W - 2*M, Inches(0.75),
         text=title, size=30, bold=True, color=color)
    bar = slide.shapes.add_shape(RND, M, Inches(1.1), Inches(1.6), Pt(4))
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    if subtitle:
        _txb(slide, M, Inches(1.18), W - 2*M, Inches(0.38),
             text=subtitle, size=12, color=GRAY_LT)


def _bullet_list(slide, items, x, y, w, h, size=14, color=WHITE,
                 bullet="▸ ", header=None, header_color=VIOLET_LT):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.text_frame.word_wrap = True
    tf = txb.text_frame
    first = True
    if header:
        p = tf.paragraphs[0]
        first = False
        run = p.add_run()
        run.text = header
        run.font.size = Pt(size + 2)
        run.font.bold = True
        run.font.color.rgb = header_color
    for item in items:
        p = tf.paragraphs[0] if (first and not header) else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb


def _slide_num(slide, n: int):
    _txb(slide, W - Inches(0.8), H - Inches(0.35), Inches(0.6), Inches(0.3),
         text=str(n), size=10, color=GRAY_DK, align=PP_ALIGN.RIGHT)


def _mpl_img(fig, slide, x, y, w, h):
    """Inserta figura matplotlib como PNG transparente en el slide."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    buf.seek(0)
    slide.shapes.add_picture(buf, x, y, w, h)
    plt.close(fig)


def _accent_top(slide, x, y, w, color, radius=0.07):
    """Franja de color redondeada en la parte superior de una card."""
    bar = slide.shapes.add_shape(RND, x, y, w, Pt(5))
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# GRÁFICAS MATPLOTLIB
# ══════════════════════════════════════════════════════════════════════════════

def _chart_finances():
    """Gráfica de barras agrupadas: Ingresos / Costos / Utilidad x 3 años."""
    fig, ax = plt.subplots(figsize=(7.5, 4.2))
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")

    years = ["Año 1", "Año 2", "Año 3"]
    ingresos = [12_000,  60_000, 300_000]
    costos   = [ 8_000,  35_000, 120_000]
    utilidad = [ 4_000,  25_000, 180_000]

    x = np.arange(len(years))
    bw = 0.24

    b1 = ax.bar(x - bw, ingresos, bw, label="Ingresos",  color=_EME,  alpha=0.92, zorder=3)
    b2 = ax.bar(x,       costos,  bw, label="Costos",    color=_RED,  alpha=0.85, zorder=3)
    b3 = ax.bar(x + bw,  utilidad,bw, label="Utilidad",  color=_VIOL, alpha=0.92, zorder=3)

    ax.set_xticks(x)
    ax.set_xticklabels(years, color=_WHT, fontsize=13, fontweight="bold")
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"${v/1000:.0f}K"))
    ax.tick_params(axis="y", colors=_GRY, labelsize=10)
    ax.tick_params(axis="x", colors=_GRY)
    ax.spines["bottom"].set_color(_GRY)
    ax.spines["left"].set_color(_GRY)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.grid(True, color=_GRY, alpha=0.2, zorder=0)
    ax.set_axisbelow(True)

    for bar, col in [(b1, _EME), (b3, _VIOL)]:
        for b in bar:
            h = b.get_height()
            ax.text(b.get_x() + b.get_width() / 2, h + 4500,
                    f"${h/1000:.0f}K", ha="center", va="bottom",
                    color=col, fontsize=8, fontweight="bold")

    legend = ax.legend(
        facecolor="#1E293B", edgecolor=_GRY, labelcolor=_WHT,
        fontsize=10, loc="upper left", framealpha=0.8
    )
    plt.tight_layout(pad=0.3)
    return fig


def _chart_market():
    """Círculos concéntricos TAM / SAM / SOM."""
    fig, ax = plt.subplots(figsize=(5, 5))
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    ax.set_aspect("equal")

    # Radios
    tam = plt.Circle((0, 0), 2.9, color="#1E293B", linewidth=2, ec=_GRY, linestyle="--")
    sam = plt.Circle((0, 0), 1.9, color="#2D1B5E", linewidth=2, ec=_VIOL)
    som = plt.Circle((0, 0), 0.85, color=_VIO, linewidth=0)

    ax.add_patch(tam)
    ax.add_patch(sam)
    ax.add_patch(som)

    ax.text(0,  3.2,  "TAM  $6.78B",    ha="center", va="center", color=_GRY,  fontsize=11)
    ax.text(0,  2.1,  "SAM  $800M",     ha="center", va="center", color=_VIOL, fontsize=11, fontweight="bold")
    ax.text(0,  0.0,  "SOM",            ha="center", va="center", color=_WHT,  fontsize=12, fontweight="bold")
    ax.text(0, -0.45, "$45M",           ha="center", va="center", color=_WHT,  fontsize=11, fontweight="bold")

    ax.set_xlim(-4, 4)
    ax.set_ylim(-4, 4)
    ax.axis("off")
    plt.tight_layout(pad=0)
    return fig


def _chart_investment_pie():
    """Donut de distribución de inversión $4,100."""
    fig, ax = plt.subplots(figsize=(4, 4))
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")

    sizes  = [1200, 2000, 500, 300, 100]
    labels = ["Laptop\n$1,200", "Capital\ntrabajo\n$2,000",
              "Legal\n$500", "Branding\n$300", "Hosting\n$100"]
    colors = [_VIOL, _EME, _AMB, _VIO, _GRY]
    explode = (0, 0.06, 0, 0, 0)

    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, colors=colors, explode=explode,
        autopct="%1.0f%%", startangle=90,
        wedgeprops={"width": 0.55, "edgecolor": _BG, "linewidth": 2},
        textprops={"color": _WHT, "fontsize": 8}
    )
    for at in autotexts:
        at.set_color(_BG)
        at.set_fontweight("bold")
        at.set_fontsize(8)

    ax.text(0, 0, "$4,100", ha="center", va="center",
            color=_WHT, fontsize=14, fontweight="bold")

    ax.axis("equal")
    plt.tight_layout(pad=0)
    return fig


# ══════════════════════════════════════════════════════════════════════════════
# 15 SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    s = _new_slide(prs)

    # Panel izquierdo con degradado visual
    _rect(s, Inches(0), Inches(0), Inches(5.0), H,
          fill=RGBColor(0x16, 0x0D, 0x38), rounded=False)

    # Líneas decorativas verticales
    for i, col in enumerate([VIOLET, EMERALD, VIOLET_LT]):
        bar = s.shapes.add_shape(FLAT, Inches(5.0), Inches(i * 2.5), Pt(3), Inches(2.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background()

    # Nombre
    _txb(s, Inches(0.5), Inches(1.2), Inches(4.2), Inches(1.0),
         "DocuMind AI", size=42, bold=True, color=WHITE)
    # Slogan
    _txb(s, Inches(0.5), Inches(2.25), Inches(4.2), Inches(0.55),
         '"El cerebro digital de tu empresa"',
         size=13, italic=True, color=VIOLET_LT)
    # Badge caso innovación
    _badge(s, Inches(0.5), Inches(2.9), Inches(2.2), Inches(0.32),
           "★  Caso de Innovación", bg=VIOLET, size=10)
    # Info grupo
    _txb(s, Inches(0.5), Inches(3.38), Inches(4.2), Inches(0.3),
         "Grupo 1GS241  ·  UTP  ·  Innovación y Emprendimiento",
         size=10, color=GRAY_LT)
    _txb(s, Inches(0.5), Inches(3.7), Inches(4.2), Inches(0.3),
         "Prof. Melvin Falcón  ·  29/05/2026", size=10, color=GRAY_LT)

    # Equipo
    team = [
        ("Kelvin He",           "8-999-1950",   "CEO"),
        ("Roy Barrera",         "8-1022-2121",  "Backend Lead"),
        ("Einer Mosquera",      "8-924-1880",   "Frontend Lead"),
        ("Nicolás Athanasidis", "8-1001-974",   "QA / Docs"),
    ]
    roles_col = [VIOLET_LT, EMERALD, AMBER, GRAY_LT]
    for i, ((name, cedula, role), col) in enumerate(zip(team, roles_col)):
        yy = Inches(4.3) + i * Inches(0.62)
        _rect(s, Inches(0.5), yy, Inches(4.2), Inches(0.55), fill=BG_CARD2, radius=0.08)
        _txb(s, Inches(0.65), yy + Pt(4), Inches(2.0), Inches(0.42),
             name, size=10, bold=True, color=WHITE)
        _txb(s, Inches(2.7), yy + Pt(4), Inches(1.1), Inches(0.42),
             cedula, size=9, color=GRAY_LT)
        _badge(s, Inches(3.85), yy + Pt(6), Inches(0.75), Inches(0.26),
               role, bg=col, size=7)

    # Panel derecho — mockup chat
    _rect(s, Inches(5.3), Inches(0.7), Inches(7.7), Inches(6.4), fill=BG_CARD, radius=0.05)
    _txb(s, Inches(5.55), Inches(0.9), Inches(5.0), Inches(0.38),
         "◉  DocuMind AI — Chat con Constructora Nova S.A.",
         size=11, bold=True, color=VIOLET_LT)

    # Separador
    sep = s.shapes.add_shape(FLAT, Inches(5.3), Inches(1.35), Inches(7.7), Pt(1))
    sep.fill.solid(); sep.fill.fore_color.rgb = BG_CARD2; sep.line.fill.background()

    msgs = [
        ("👤  ¿Cuántos días de vacaciones con 2 años de trabajo?", GRAY_LT,  PP_ALIGN.RIGHT),
        ("🤖  Según el Manual de RRHH (p.4), un empleado con 2 años", EMERALD, PP_ALIGN.LEFT),
        ("      tiene 17 días hábiles de vacaciones anuales.",          EMERALD, PP_ALIGN.LEFT),
        ("      Fuente: Manual_RRHH_Nova.pdf, página 4.",               GRAY_LT, PP_ALIGN.LEFT),
        ("",                                                            WHITE,   PP_ALIGN.LEFT),
        ("👤  ¿Quién debe aprobar una compra de $3,000?",              GRAY_LT, PP_ALIGN.RIGHT),
        ("🤖  El Gerente General junto con el Gerente de Finanzas.",   EMERALD, PP_ALIGN.LEFT),
        ("      Fuente: Politica_Compras_Proveedores.pdf, página 4.",  GRAY_LT, PP_ALIGN.LEFT),
        ("",                                                            WHITE,   PP_ALIGN.LEFT),
        ("👤  ¿Cuál es el plazo del contrato de arrendamiento?",       GRAY_LT, PP_ALIGN.RIGHT),
        ("🤖  El contrato tiene un plazo de 24 meses renovables.",     EMERALD, PP_ALIGN.LEFT),
        ("      Fuente: Contrato_Arrendamiento_Oficina.pdf, página 2.",GRAY_LT, PP_ALIGN.LEFT),
    ]
    for i, (msg, col, align) in enumerate(msgs):
        _txb(s, Inches(5.5), Inches(1.45) + i * Inches(0.38), Inches(7.3), Inches(0.36),
             msg, size=10, color=col, align=align)

    _txb(s, Inches(5.55), H - Inches(0.6), Inches(7.0), Inches(0.3),
         "7 docs indexados  ·  1,847 chunks  ·  Groq llama-3.1-8b  ·  ~0.8 s/respuesta",
         size=9, color=GRAY_DK)
    _slide_num(s, 1)


def slide_02_problem(prs):
    s = _new_slide(prs)
    _slide_title(s, "El Problema",
                 "Las empresas pierden miles de horas buscando información en sus propios documentos")

    stats = [
        ("2.5 h",    "perdidas por empleado\ncada semana buscando\ndocumentos internos",     VIOLET_LT),
        ("6,000 h",  "anuales de productividad\nperdida en una empresa\nde 50 empleados",     EMERALD),
        ("$75,000",  "costo anual estimado\nen tiempo no productivo\n(empresa 50 empleados)", AMBER),
        ("95%",      "de MIPYMES panameñas\nsin herramientas de gestión\ndocumental con IA",  RED_SOFT),
    ]
    cw = Inches(3.1)
    for i, (val, desc, col) in enumerate(stats):
        cx = M + i * (cw + Inches(0.08))
        _rect(s, cx, Inches(1.55), cw, Inches(2.5), fill=BG_CARD, radius=0.08)
        _accent_top(s, cx, Inches(1.55), cw, col)
        _txb(s, cx + Inches(0.1), Inches(1.75), cw - Inches(0.2), Inches(0.9),
             val, size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
        _txb(s, cx + Inches(0.1), Inches(2.7), cw - Inches(0.2), Inches(1.1),
             desc, size=11, color=GRAY_LT, align=PP_ALIGN.CENTER)

    _bullet_list(s, [
        "Carpetas compartidas caóticas: nadie sabe dónde está nada.",
        "Empleados preguntan a compañeros → interrumpen la productividad.",
        "Versiones obsoletas de contratos generan errores y conflictos legales.",
        "Búsquedas manuales en cientos de PDFs consumen horas cada semana.",
        "Conocimiento crítico se pierde cuando un empleado clave se va.",
    ], M, Inches(4.3), W - 2*M, Inches(2.5),
    size=15, header="¿Cómo se ve hoy en una empresa panameña?")

    _txb(s, M, H - Inches(0.45), W - 2*M, Inches(0.28),
         "Fuentes: AMPYME 2026  ·  La Prensa Panamá  ·  Estimaciones internas DocuMind AI",
         size=8, color=GRAY_DK)
    _slide_num(s, 2)


def slide_03_solution(prs):
    s = _new_slide(prs)
    _slide_title(s, "La Solución — DocuMind AI",
                 "Convierte documentos estáticos en conocimiento accionable mediante IA")

    # Quote card
    _rect(s, M, Inches(1.6), W - 2*M, Inches(0.85), fill=BG_CARD2, radius=0.06)
    _txb(s, Inches(0.75), Inches(1.72), W - Inches(1.5), Inches(0.65),
         "ChatGPT responde con conocimiento general de internet.  "
         "DocuMind responde con los documentos internos de tu empresa, "
         "citando archivo, página y párrafo exactos.",
         size=14, italic=True, color=WHITE)

    features = [
        ("🔍", "DocuMind Search",    VIOLET,
         "Búsqueda semántica\nque entiende el significado,\nno solo palabras clave.\nEncuentra sinónimos y contexto."),
        ("💬", "DocuMind Chat",     EMERALD,
         "Chatbot RAG que responde\nen lenguaje natural citando\narchivo, página y párrafo\nexactos. Siempre verificable."),
        ("📷", "OCR Integrado",     AMBER,
         "Procesa PDFs escaneados\ne imágenes JPG/PNG con\nTesseract. Muestra % de\nconfianza del OCR."),
        ("📊", "Analytics",        VIOLET_LT,
         "Dashboard con métricas\nen tiempo real: documentos,\nchunks, consultas,\nactividad 7 días."),
        ("💳", "Free & Paid Plans", EMERALD,
         "Plan Free de entrada.\nUpgrade en 2 clics vía\nStripe. Plan activo al\ninstante sin ngrok."),
    ]
    cw = Inches(2.48)
    for i, (icon, title, col, desc) in enumerate(features):
        cx = M + i * (cw + Inches(0.08))
        _rect(s, cx, Inches(2.68), cw, Inches(4.45), fill=BG_CARD, radius=0.07)
        _accent_top(s, cx, Inches(2.68), cw, col)
        _txb(s, cx + Inches(0.1), Inches(2.82), cw - Inches(0.2), Inches(0.55),
             icon, size=24, align=PP_ALIGN.CENTER)
        _txb(s, cx + Inches(0.1), Inches(3.42), cw - Inches(0.2), Inches(0.52),
             title, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        _txb(s, cx + Inches(0.12), Inches(3.98), cw - Inches(0.24), Inches(2.8),
             desc, size=10.5, color=GRAY_LT, align=PP_ALIGN.CENTER)

    _slide_num(s, 3)


def slide_04_demo(prs):
    s = _new_slide(prs)
    _slide_title(s, "Demo — Constructora Nova S.A.",
                 "7 documentos reales · 12 preguntas · respuestas con cita exacta en < 1 segundo")

    # Panel izquierdo — flujo demo
    _rect(s, M, Inches(1.55), Inches(4.6), Inches(5.55), fill=BG_CARD, radius=0.06)
    _txb(s, Inches(0.68), Inches(1.72), Inches(4.2), Inches(0.38),
         "⏱  Flujo de Demo — 3 minutos", size=13, bold=True, color=VIOLET_LT)

    pasos = [
        ("1", "Landing futurista animada",        "15 s",  VIOLET),
        ("2", "Login con Clerk (cuenta real)",     "10 s",  VIOLET_LT),
        ("3", "Dashboard con 7 docs indexados",   "20 s",  EMERALD),
        ("4", "Subir imagen → OCR automático",    "30 s",  AMBER),
        ("5", "Chat con 3 preguntas y citas",     "90 s",  EMERALD),
        ("6", "Búsqueda semántica con score",     "20 s",  VIOLET_LT),
        ("7", "Free → Pago en vivo (Stripe)",     "15 s",  AMBER),
    ]
    for i, (num, step, time, col) in enumerate(pasos):
        yy = Inches(2.22) + i * Inches(0.62)
        _badge(s, Inches(0.68), yy + Pt(2), Inches(0.34), Inches(0.34), num, bg=col, size=10)
        _txb(s, Inches(1.12), yy + Pt(4), Inches(2.85), Inches(0.34), step, size=11, color=WHITE)
        _txb(s, Inches(4.1), yy + Pt(4), Inches(0.8), Inches(0.34),
             time, size=10, color=GRAY_LT, align=PP_ALIGN.RIGHT)

    # Panel derecho — Q&A
    _rect(s, Inches(5.4), Inches(1.55), Inches(7.5), Inches(5.55), fill=BG_CARD, radius=0.06)
    _txb(s, Inches(5.6), Inches(1.72), Inches(7.0), Inches(0.38),
         "❓  Preguntas del Demo (selección de las 12)", size=13, bold=True, color=EMERALD)

    qa = [
        ("¿Cuántos días de vacaciones con 2 años?",
         "▸  17 días hábiles  [Manual RRHH p.4]"),
        ("¿Preaviso para renunciar?",
         "▸  30 días calendario  [Contrato Laboral p.3]"),
        ("¿Plazo del contrato de arrendamiento?",
         "▸  24 meses renovables  [Arrendamiento p.2]"),
        ("¿Quién aprueba compras de $3,000?",
         "▸  Gerente General + Finanzas  [Compras p.4]"),
        ("¿Cobertura del seguro médico en emergencias?",
         "▸  100% sin deducible  [Contrato Laboral p.4]"),
        ("¿Qué dice el acta sobre el presupuesto?",
         "▸  $3,200,000 aprobado (OCR)  [Acta p.1]"),
    ]
    for i, (q, a) in enumerate(qa):
        yy = Inches(2.22) + i * Inches(0.75)
        _txb(s, Inches(5.6), yy, Inches(7.1), Inches(0.33), q, size=11, color=WHITE)
        _txb(s, Inches(5.6), yy + Inches(0.32), Inches(7.1), Inches(0.33),
             a, size=10, italic=True, color=EMERALD)

    _badge(s, Inches(5.6), H - Inches(0.65), Inches(2.5), Inches(0.28),
           "Groq llama-3.1-8b · ~0.8 s/respuesta", bg=BG_CARD2, size=9)
    _badge(s, Inches(8.3), H - Inches(0.65), Inches(2.3), Inches(0.28),
           "Fallback extractivo sin LLM", bg=BG_CARD2, size=9)
    _slide_num(s, 4)


def slide_05_tech(prs):
    s = _new_slide(prs)
    _slide_title(s, "Tecnología",
                 "Stack moderno y open-source · Sin Docker · APIs gratuitas · Fallback sin internet")

    layers = [
        ("Frontend",         VIOLET_LT, [
            "Next.js 14", "React 18", "TypeScript", "Tailwind v4", "shadcn/ui", "Framer Motion", "Recharts",
        ]),
        ("Auth & Pagos",     VIOLET, [
            "Clerk 7", "Stripe 22", "OAuth + Email", "Stripe Checkout", "Sin webhooks",
        ]),
        ("Backend API",      EMERALD, [
            "FastAPI 0.138", "Python 3.11", "Uvicorn --reload", "4 endpoints REST", "Pydantic v2",
        ]),
        ("LLM — Groq",       AMBER, [
            "llama-3.1-8b-instant", "30 req/min gratis", "~0.8 s respuesta", "Fallback extractivo",
        ]),
        ("Embeddings",       VIOLET_LT, [
            "all-MiniLM-L6-v2", "ChromaDB ONNX", "22 MB", "Sin PyTorch", "Búsqueda local offline",
        ]),
        ("OCR",              EMERALD, [
            "Tesseract OCR", "pdf2image", "Poppler", "Pillow", "Grayscale + binarización", "Confianza %",
        ]),
        ("Almacenamiento",   GRAY_LT, [
            "ChromaDB 1.5 (vectores)", "SQLite (metadata)", "Sistema de archivos (uploads)",
        ]),
    ]

    for i, (layer, col, chips) in enumerate(layers):
        yy = Inches(1.55) + i * Inches(0.715)
        # Etiqueta de capa
        _rect(s, M, yy, Inches(2.6), Inches(0.62), fill=col, radius=0.08)
        lbl = s.shapes[-1]
        tf = lbl.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = layer
        run.font.size = Pt(11.5)
        run.font.bold = True
        run.font.color.rgb = BG_DARK

        # Chips de tecnologías
        chip_x = Inches(3.25)
        for chip in chips:
            cw = Inches(len(chip) * 0.09 + 0.5)
            if chip_x + cw > W - M:
                break
            _badge(s, chip_x, yy + Inches(0.14), cw, Inches(0.34),
                   chip, bg=BG_CARD2, tc=WHITE, size=10)
            chip_x += cw + Inches(0.12)

    _slide_num(s, 5)


def slide_06_product(prs):
    s = _new_slide(prs)
    _slide_title(s, "El Producto",
                 "Sistema Operativo del Conocimiento Empresarial")

    features = [
        ("🔍", "Búsqueda Semántica",  VIOLET,
         "Entiende el significado real.\n'Baja por maternidad' encuentra\n'permiso por maternidad'.\nSin coincidencia exacta de palabras."),
        ("💬", "Chat con Citas",      EMERALD,
         "Responde en lenguaje natural\ncitando archivo, página y\npárrafo exactos.\nSiempre verificable, nunca inventa."),
        ("📷", "OCR Avanzado",        AMBER,
         "PDFs escaneados e imágenes\nJPG/PNG con Tesseract.\nPreprocesamiento automático.\nConfianza % visible en UI."),
        ("📊", "Dashboard Analytics", VIOLET_LT,
         "Métricas en tiempo real:\ndocumentos, chunks, consultas,\nalmacenamiento y gráfica\nde actividad (7 días)."),
        ("💳", "Gating Free/Pago",   EMERALD,
         "Free: 3 docs, sin OCR.\nUpgrade en 2 clics vía Stripe.\nPlan sube al instante.\nSin ngrok ni webhooks."),
        ("🛡️", "Fallback Extractivo", GRAY_LT,
         "Sin Groq → muestra los chunks\nmás relevantes directamente.\n100% uptime. La app nunca\nse cae por falta de LLM."),
    ]
    cw = Inches(4.08)
    ch = Inches(2.65)
    for i, (icon, title, col, desc) in enumerate(features):
        col_i = i % 3
        row_i = i // 3
        cx = M + col_i * (cw + Inches(0.17))
        cy = Inches(1.58) + row_i * (ch + Inches(0.15))
        _rect(s, cx, cy, cw, ch, fill=BG_CARD, radius=0.07)
        _accent_top(s, cx, cy, cw, col)
        _txb(s, cx + Inches(0.15), cy + Inches(0.14), Inches(0.55), Inches(0.5),
             icon, size=22)
        _txb(s, cx + Inches(0.75), cy + Inches(0.18), cw - Inches(0.9), Inches(0.45),
             title, size=13, bold=True, color=col)
        _txb(s, cx + Inches(0.15), cy + Inches(0.65), cw - Inches(0.3), Inches(1.85),
             desc, size=10.5, color=GRAY_LT)

    _slide_num(s, 6)


def slide_07_market(prs):
    s = _new_slide(prs)
    _slide_title(s, "Mercado Objetivo",
                 "Comenzando en Panamá · Expandiendo por Latinoamérica · TAM global $6.78B")

    # Gráfica de círculos concéntricos
    fig = _chart_market()
    _mpl_img(fig, s, Inches(7.6), Inches(1.4), Inches(5.5), Inches(5.5))

    # Segmentación
    _bullet_list(s, [
        "Despachos legales (bufetes, notarías, CSJ)",
        "Constructoras (gestión de contratos y RRHH)",
        "Consultoras (auditoría, recursos humanos, finanzas)",
        "Universidades (reglamentos, actas, administrativo)",
        "Bancos y aseguradoras (compliance y riesgos)",
    ], M, Inches(1.65), Inches(6.8), Inches(3.0),
    size=14, header="Fase 1 — Panamá (Meses 1-12)")

    _bullet_list(s, [
        "Año 2: Costa Rica + Colombia",
        "Año 3+: México, Perú, Chile, Ecuador",
        "Buyer persona: Gerente PYME con > 50 docs/mes",
    ], M, Inches(4.85), Inches(6.8), Inches(1.8),
    size=13, header="Expansión regional")

    # KPI cards
    kpis = [("TAM", "$6.78B", GRAY_LT), ("SAM", "$800M", VIOLET_LT), ("SOM", "$45M", EMERALD)]
    for i, (label, val, col) in enumerate(kpis):
        cx = M + i * Inches(2.3)
        _rect(s, cx, H - Inches(0.85), Inches(2.1), Inches(0.62), fill=BG_CARD2, radius=0.1)
        _txb(s, cx + Inches(0.1), H - Inches(0.82), Inches(0.7), Inches(0.52),
             label, size=11, bold=True, color=col)
        _txb(s, cx + Inches(0.8), H - Inches(0.82), Inches(1.2), Inches(0.52),
             val, size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    _slide_num(s, 7)


def slide_08_competition(prs):
    s = _new_slide(prs)
    _slide_title(s, "Competencia",
                 "Ningún competidor muestra 'página 7, párrafo 3' — esa es nuestra ventaja diferencial")

    headers = ["Característica", "DocuMind AI", "ChatGPT", "Copilot", "Glean", "Notion AI"]
    rows = [
        ["Citas exactas (archivo + página)", "✅ SÍ",     "❌ NO",       "⚠️ PARCIAL",  "❌ NO",       "❌ NO"],
        ["Multi-documento simultáneo",        "✅ SÍ",     "❌ NO",       "✅ SÍ",        "✅ SÍ",        "❌ NO"],
        ["Español latinoamericano nativo",    "✅ NATIVO", "⚠️ GENÉRICO", "⚠️ GENÉRICO", "⚠️ GENÉRICO", "⚠️ GENÉRICO"],
        ["OCR (imágenes y PDFs escaneados)", "✅ SÍ",     "❌ NO",       "❌ NO",        "❌ NO",        "❌ NO"],
        ["On-premise opcional",               "✅ SÍ",     "❌ NO",       "❌ NO",        "❌ NO",        "❌ NO"],
        ["Precio accesible para MIPYMES",     "✅ $49/mes","⚠️ $30/user", "⚠️ $30/user", "❌ $500+/mes", "⚠️ $10/user"],
        ["Stack open source (sin lock-in)",   "✅ SÍ",     "❌ NO",       "❌ NO",        "❌ NO",        "❌ NO"],
    ]

    col_widths = [Inches(3.5), Inches(1.9), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)]
    tbl_x, tbl_y = M, Inches(1.55)
    tbl_w = sum(col_widths)
    tbl = s.shapes.add_table(len(rows) + 1, len(headers),
                              tbl_x, tbl_y, tbl_w, Inches(5.1)).table

    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    def _cell(cell, text, bg, tc=WHITE, sz=10, bold=False):
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz); run.font.bold = bold
        run.font.color.rgb = tc

    for j, h in enumerate(headers):
        _cell(tbl.rows[0].cells[j], h,
              bg=VIOLET if j == 1 else BG_CARD2, bold=True, sz=10)

    for i, row in enumerate(rows):
        bg_row = BG_DARK if i % 2 == 0 else BG_CARD
        for j, val in enumerate(row):
            if j == 0:
                _cell(tbl.rows[i+1].cells[j], val, bg=bg_row, tc=GRAY_LT, sz=10)
                tbl.rows[i+1].cells[j].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif j == 1:
                col_v = EMERALD if "✅" in val else (RED_SOFT if "❌" in val else AMBER)
                _cell(tbl.rows[i+1].cells[j], val, bg=BG_CARD, tc=col_v, bold=True, sz=10)
            else:
                col_v = EMERALD if "✅" in val else (RED_SOFT if "❌" in val else AMBER)
                _cell(tbl.rows[i+1].cells[j], val, bg=bg_row, tc=col_v, sz=10)

    _slide_num(s, 8)


def slide_09_business(prs):
    s = _new_slide(prs)
    _slide_title(s, "Modelo de Negocio",
                 "SaaS con freemium · Upgrade en 2 clics vía Stripe · Sin contratos anuales obligatorios")

    plans = [
        ("Free",       "$0/mes",    BG_CARD,  GRAY_LT,   False,
         ["3 documentos máx.", "20 consultas/mes", "Sin OCR (imágenes)", "1 usuario", "— soporte community"]),
        ("Starter",    "$49/mes",   VIOLET,   WHITE,     True,
         ["Docs ilimitados", "OCR activo", "Chat RAG + búsqueda", "5 usuarios", "Soporte email"]),
        ("Business",   "$199/mes",  BG_CARD,  GRAY_LT,   False,
         ["50 usuarios", "100 GB storage", "Multi-workspace", "API access", "Soporte prioritario"]),
        ("Enterprise", "$999/mes",  BG_CARD,  GRAY_LT,   False,
         ["Usuarios ilimitados", "On-premise", "SSO/SAML", "CSM dedicado", "SLA 99.9%"]),
    ]
    cw = Inches(3.1)
    for i, (name, price, fill, tcol, recommended, feats) in enumerate(plans):
        cx = M + i * (cw + Inches(0.1))
        _rect(s, cx, Inches(1.55), cw, Inches(5.55), fill=fill, radius=0.07)
        if recommended:
            _badge(s, cx + Inches(0.55), Inches(1.55), Inches(2.0), Inches(0.3),
                   "⭐  RECOMENDADO", bg=EMERALD, size=8)
        name_col = EMERALD if recommended else VIOLET_LT
        _txb(s, cx + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(0.52),
             name, size=20, bold=True, color=name_col, align=PP_ALIGN.CENTER)
        _txb(s, cx + Inches(0.1), Inches(2.52), cw - Inches(0.2), Inches(0.6),
             price, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Separador
        sep = s.shapes.add_shape(FLAT, cx + Inches(0.2), Inches(3.18), cw - Inches(0.4), Pt(1))
        sep.fill.solid(); sep.fill.fore_color.rgb = BG_CARD2 if not recommended else VIOLET_LT
        sep.line.fill.background()

        for j, feat in enumerate(feats):
            _txb(s, cx + Inches(0.2), Inches(3.3) + j * Inches(0.43),
                 cw - Inches(0.4), Inches(0.38),
                 f"▸  {feat}", size=10.5, color=tcol)

    # Stripe flow
    _rect(s, M, H - Inches(0.75), W - 2*M, Inches(0.5), fill=BG_CARD2, radius=0.08)
    _txb(s, Inches(0.7), H - Inches(0.7), W - Inches(1.4), Inches(0.4),
         "Flujo de upgrade:  UpgradeDialog  →  POST /api/stripe/checkout  →  Stripe Checkout  →  "
         "POST /api/stripe/confirm  →  publicMetadata.plan='starter'  →  UI refleja Starter al instante",
         size=9.5, color=GRAY_LT)

    _slide_num(s, 9)


def slide_10_traction(prs):
    s = _new_slide(prs)
    _slide_title(s, "Tracción Inicial",
                 "Validación antes de inversión — lean startup · enfoque en señales reales del mercado")

    metrics = [
        ("50+",  "Empresas en\nwaitlist activa",   VIOLET_LT),
        ("10",   "Pilotos\ncomprometidos",          EMERALD),
        ("< 2s", "Tiempo de respuesta\npromedio",   AMBER),
        ("$0",   "Inversión en\nadquisición hoy",   GRAY_LT),
    ]
    cw = Inches(3.08)
    for i, (val, label, col) in enumerate(metrics):
        cx = M + i * (cw + Inches(0.08))
        _rect(s, cx, Inches(1.58), cw, Inches(2.3), fill=BG_CARD, radius=0.08)
        _accent_top(s, cx, Inches(1.58), cw, col)
        _txb(s, cx + Inches(0.1), Inches(1.75), cw - Inches(0.2), Inches(0.95),
             val, size=46, bold=True, color=col, align=PP_ALIGN.CENTER)
        _txb(s, cx + Inches(0.1), Inches(2.72), cw - Inches(0.2), Inches(0.75),
             label, size=12, color=GRAY_LT, align=PP_ALIGN.CENTER)

    # Testimonios
    testimonios = [
        ("🏗️  Nova Constructora S.A.",
         '"Ahora cualquier empleado sabe sus cláusulas de\ncontrato en segundos, sin llamarme a mí."',
         "Carlos Méndez · Gerente General"),
        ("⚖️  Bufete Méndez & Asociados",
         '"Encontrar jurisprudencia en 15,000 PDFs\nsolía tomar 3 horas. Ahora: 8 segundos."',
         "Dr. Luis Méndez · Socio fundador"),
        ("🎓  Universidad del Pacífico",
         '"Los estudiantes consultan reglamentos sin\ncolapsar la línea de secretaría."',
         "Lic. Sofía Arias · Directora Académica"),
    ]
    cw_t = Inches(4.1)
    for i, (company, quote, person) in enumerate(testimonios):
        cx = M + i * (cw_t + Inches(0.18))
        _rect(s, cx, Inches(4.1), cw_t, Inches(2.9), fill=BG_CARD, radius=0.07)
        _txb(s, cx + Inches(0.15), Inches(4.25), cw_t - Inches(0.3), Inches(0.35),
             company, size=11, bold=True, color=VIOLET_LT)
        _txb(s, cx + Inches(0.15), Inches(4.65), cw_t - Inches(0.3), Inches(1.4),
             quote, size=11, italic=True, color=WHITE)
        _txb(s, cx + Inches(0.15), Inches(6.1), cw_t - Inches(0.3), Inches(0.3),
             f"— {person}", size=9, color=GRAY_LT)

    _slide_num(s, 10)


def slide_11_financials(prs):
    s = _new_slide(prs)
    _slide_title(s, "Finanzas — Proyección 3 Años",
                 "Break-even en el mes 10 · Margen 33% → 60% · LTV/CAC = 28.8×")

    # Gráfica matplotlib
    fig = _chart_finances()
    _mpl_img(fig, s, Inches(0.4), Inches(1.45), Inches(7.8), Inches(5.0))

    # KPIs derecha
    kpis = [
        ("Clientes — Fin Año 1",  "20",         VIOLET_LT),
        ("MRR — Año 1",           "$2,040",      EMERALD),
        ("ARR — Año 1",           "$12,000",     EMERALD),
        ("Clientes — Fin Año 3",  "340",         VIOLET_LT),
        ("MRR — Año 3",           "$34,680",     EMERALD),
        ("ARR — Año 3",           "$300,000",    EMERALD),
        ("Margen Año 1",          "33%",         AMBER),
        ("Margen Año 3",          "60%",         AMBER),
        ("CAC",                   "$50/cliente", GRAY_LT),
        ("LTV / CAC",             "28.8 ×",      AMBER),
        ("Break-even",            "Mes 10",      VIOLET_LT),
    ]
    _txb(s, Inches(8.5), Inches(1.48), Inches(4.5), Inches(0.42),
         "KPIs Clave", size=13, bold=True, color=VIOLET_LT)
    for i, (key, val, col) in enumerate(kpis):
        yy = Inches(1.97) + i * Inches(0.42)
        bg = BG_CARD if i % 2 == 0 else BG_CARD2
        _rect(s, Inches(8.5), yy, Inches(4.5), Inches(0.38), fill=bg, radius=0.06)
        _txb(s, Inches(8.65), yy + Pt(3), Inches(2.7), Inches(0.3),
             key, size=9.5, color=GRAY_LT)
        _txb(s, Inches(11.1), yy + Pt(3), Inches(1.7), Inches(0.3),
             val, size=10, bold=True, color=col, align=PP_ALIGN.RIGHT)

    _slide_num(s, 11)


def slide_12_team(prs):
    s = _new_slide(prs)
    _slide_title(s, "El Equipo", "Grupo 1GS241 · UTP · Innovación y Emprendimiento · 29/05/2026")

    members = [
        ("Kelvin He",           "8-999-1950",   "CEO",           VIOLET,
         "Visión del producto\nPitch a inversionistas\nEstrategia de negocio\nRelaciones con aliados"),
        ("Roy Barrera",         "8-1022-2121",  "Backend Lead",  EMERALD,
         "Arquitectura RAG\nFastAPI + ChromaDB\nEmbeddings ONNX\nOCR + Tesseract"),
        ("Einer Mosquera",      "8-924-1880",   "Frontend Lead", AMBER,
         "Next.js 14 + React\nshadcn/ui + Tailwind\nLanding futurista\nDashboard + Chat UI"),
        ("Nicolás Athanasidis", "8-1001-974",   "QA / Docs",     VIOLET_LT,
         "Pruebas funcionales\nDocumentación técnica\nSoporte en el pitch\nControl de calidad"),
    ]

    # CEO card centrada arriba
    ceo = members[0]
    ceo_x = Inches(4.8)
    _rect(s, ceo_x, Inches(1.55), Inches(3.7), Inches(2.8), fill=BG_CARD, radius=0.08)
    _accent_top(s, ceo_x, Inches(1.55), Inches(3.7), ceo[3])
    _txb(s, ceo_x + Inches(0.15), Inches(1.75), Inches(3.4), Inches(0.52),
         ceo[0], size=18, bold=True, color=ceo[3], align=PP_ALIGN.CENTER)
    _txb(s, ceo_x + Inches(0.15), Inches(2.27), Inches(3.4), Inches(0.28),
         ceo[1], size=10, color=GRAY_LT, align=PP_ALIGN.CENTER)
    _badge(s, ceo_x + Inches(1.1), Inches(2.6), Inches(1.5), Inches(0.28),
           ceo[2], bg=ceo[3], size=10)
    _txb(s, ceo_x + Inches(0.15), Inches(2.97), Inches(3.4), Inches(1.2),
         ceo[4], size=10, color=GRAY_LT, align=PP_ALIGN.CENTER)

    # Líneas del organigrama
    def _line(slide, x, y, w, h, col=GRAY_DK):
        ln = slide.shapes.add_shape(FLAT, x, y, w, h)
        ln.fill.solid(); ln.fill.fore_color.rgb = col; ln.line.fill.background()

    _line(s, Inches(6.65), Inches(4.35), Pt(2), Inches(0.45))  # vertical CEO
    _line(s, Inches(2.1),  Inches(4.8),  Inches(9.2), Pt(2))   # horizontal

    # 3 miembros
    for i, m in enumerate(members[1:]):
        cx = Inches(0.55) + i * Inches(4.35)
        _line(s, cx + Inches(1.78), Inches(4.8), Pt(2), Inches(0.38))  # vertical bajada
        _rect(s, cx, Inches(5.18), Inches(3.6), Inches(1.95), fill=BG_CARD, radius=0.08)
        _accent_top(s, cx, Inches(5.18), Inches(3.6), m[3])
        _txb(s, cx + Inches(0.1), Inches(5.35), Inches(3.4), Inches(0.45),
             m[0], size=13, bold=True, color=m[3], align=PP_ALIGN.CENTER)
        _txb(s, cx + Inches(0.1), Inches(5.8), Inches(3.4), Inches(0.25),
             m[1], size=9, color=GRAY_LT, align=PP_ALIGN.CENTER)
        _badge(s, cx + Inches(0.65), Inches(6.08), Inches(2.3), Inches(0.26),
               m[2], bg=m[3], size=9)
        _txb(s, cx + Inches(0.1), Inches(6.38), Inches(3.4), Inches(0.6),
             m[4].replace("\n", "  ·  "), size=8.5, color=GRAY_LT, align=PP_ALIGN.CENTER)

    _slide_num(s, 12)


def slide_13_roadmap(prs):
    s = _new_slide(prs)
    _slide_title(s, "Roadmap",
                 "De MVP universitario a infraestructura de conocimiento para Latinoamérica")

    phases = [
        ("🚀  Corto Plazo\n0 – 6 meses",   VIOLET, [
            "Lanzar MVP funcional y sin errores",
            "10 empresas piloto en Panamá",
            "Waitlist: 500 suscriptores",
            "5 pagos reales (Revenue > $0)",
            "Registrar marca DocuMind en DPI",
            "Constituir empresa (SA Panamá)",
        ]),
        ("⚡  Mediano Plazo\n6 – 18 meses", EMERALD, [
            "100 clientes pagos · MRR $10K",
            "Versión 2.0 · multi-tenancy + API",
            "Expansión: Costa Rica + Colombia",
            "Equipo de 5 personas",
            "Ronda semilla $50K – $100K",
            "3 alianzas activas (AMPYME, CC, UTP)",
        ]),
        ("🌎  Largo Plazo\n18 – 36 meses", AMBER, [
            "1,000+ clientes · ARR > $1M",
            "5 países (PA, CR, CO, MX, PE)",
            "Módulos: Compliance + Insights",
            "Integración: DGI, CSS, SAP, Slack",
            "Expansión hispanos USA (Miami)",
            "Adquisición estratégica o IPO local",
        ]),
    ]
    cw = Inches(4.1)
    for i, (phase, col, items) in enumerate(phases):
        cx = M + i * (cw + Inches(0.26))
        _rect(s, cx, Inches(1.55), cw, Inches(5.6), fill=BG_CARD, radius=0.08)
        _accent_top(s, cx, Inches(1.55), cw, col)
        _txb(s, cx + Inches(0.15), Inches(1.67), cw - Inches(0.3), Inches(0.7),
             phase, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            _txb(s, cx + Inches(0.2), Inches(2.45) + j * Inches(0.71),
                 cw - Inches(0.4), Inches(0.62),
                 f"▸  {item}", size=12, color=WHITE)

    _slide_num(s, 13)


def slide_14_investment(prs):
    s = _new_slide(prs)
    _slide_title(s, "Inversión Requerida",
                 "Bootstrapping viable · $4,100 para llegar al break-even en el mes 10")

    # Gran número + donut chart
    _txb(s, Inches(0.5), Inches(1.5), Inches(4.2), Inches(1.3),
         "$4,100", size=68, bold=True, color=VIOLET_LT, align=PP_ALIGN.CENTER)
    _txb(s, Inches(0.5), Inches(2.85), Inches(4.2), Inches(0.4),
         "Inversión inicial total (seed)", size=13, color=GRAY_LT, align=PP_ALIGN.CENTER)

    # Donut chart de distribución
    fig = _chart_investment_pie()
    _mpl_img(fig, s, Inches(0.3), Inches(3.35), Inches(4.6), Inches(3.65))

    # KPI cards
    kpis = [
        ("CAC",          "$50 / cliente",    VIOLET_LT, "Marketing orgánico + alianzas"),
        ("LTV",          "$1,440 / cliente", EMERALD,   "$120 × 12 meses × 85% retención"),
        ("LTV / CAC",    "28.8 ×",           AMBER,     "Excelente — benchmark >3 es saludable"),
        ("Break-even",   "Mes 10",           VIOLET,    "Con 2–3 clientes/mes de crecimiento"),
        ("Margen bruto", "33% → 60%",        EMERALD,   "Año 1 → Año 3 (modelo SaaS escalable)"),
    ]
    for i, (label, val, col, note) in enumerate(kpis):
        yy = Inches(1.55) + i * Inches(1.08)
        _rect(s, Inches(5.1), yy, Inches(7.8), Inches(0.95), fill=BG_CARD, radius=0.08)
        _rect(s, Inches(5.1), yy, Pt(5), Inches(0.95), fill=col, rounded=False)
        _txb(s, Inches(5.35), yy + Pt(5), Inches(2.5), Inches(0.42),
             label, size=12, bold=True, color=col)
        _txb(s, Inches(7.9), yy + Pt(5), Inches(4.8), Inches(0.42),
             val, size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        _txb(s, Inches(5.35), yy + Inches(0.56), Inches(7.4), Inches(0.3),
             note, size=9, color=GRAY_LT)

    _slide_num(s, 14)


def slide_15_cta(prs):
    s = _new_slide(prs)

    # Fondo superior degradado
    _rect(s, Inches(0), Inches(0), W, Inches(4.4),
          fill=RGBColor(0x16, 0x0D, 0x38), rounded=False)

    _txb(s, M, Inches(0.65), W - 2*M, Inches(1.2),
         "Transforma tu empresa.", size=54, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
    _txb(s, M, Inches(1.9), W - 2*M, Inches(0.65),
         "El cerebro digital de tu empresa está listo.", size=22,
         italic=True, color=VIOLET_LT, align=PP_ALIGN.CENTER)
    _txb(s, M, Inches(2.65), W - 2*M, Inches(0.55),
         '"Innovar no es solo crear algo nuevo, es crear algo que valga la pena."',
         size=13, italic=True, color=GRAY_LT, align=PP_ALIGN.CENTER)

    # 3 cards CTA
    cols_cta = [
        ("🌐", "Demo en Vivo",   VIOLET,   "http://localhost:3000",
         "Levanta el proyecto con\n2 comandos y pruébalo ahora."),
        ("📧", "Contacto",       EMERALD,  "kelvinhe04@gmail.com",
         "Primeras 50 empresas reciben\nacceso anticipado gratuito."),
        ("🎓", "UTP · 2026",     AMBER,    "Grupo 1GS241",
         "Innovación y Emprendimiento\nProf. Melvin Falcón · 29/05/2026"),
    ]
    cw = Inches(3.9)
    for i, (icon, title, col, val, note) in enumerate(cols_cta):
        cx = Inches(0.7) + i * Inches(4.3)
        _rect(s, cx, Inches(4.6), cw, Inches(2.5), fill=BG_CARD, radius=0.09)
        _accent_top(s, cx, Inches(4.6), cw, col)
        _txb(s, cx + Inches(0.15), Inches(4.75), Inches(0.6), Inches(0.5),
             icon, size=22)
        _txb(s, cx + Inches(0.8), Inches(4.8), cw - Inches(0.95), Inches(0.42),
             title, size=14, bold=True, color=col)
        _txb(s, cx + Inches(0.15), Inches(5.27), cw - Inches(0.3), Inches(0.38),
             val, size=11, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(s, cx + Inches(0.15), Inches(5.7), cw - Inches(0.3), Inches(0.85),
             note, size=10, color=GRAY_LT, align=PP_ALIGN.CENTER)

    # QR placeholder (cuadrado blanco con texto)
    _rect(s, Inches(12.25), Inches(4.55), Inches(0.85), Inches(0.85),
          fill=WHITE, rounded=True, radius=0.1)
    _txb(s, Inches(12.25), Inches(4.55), Inches(0.85), Inches(0.85),
         "QR\nDemo", size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    _slide_num(s, 15)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def build_deck() -> None:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_demo(prs)
    slide_05_tech(prs)
    slide_06_product(prs)
    slide_07_market(prs)
    slide_08_competition(prs)
    slide_09_business(prs)
    slide_10_traction(prs)
    slide_11_financials(prs)
    slide_12_team(prs)
    slide_13_roadmap(prs)
    slide_14_investment(prs)
    slide_15_cta(prs)

    prs.save(OUTPUT)
    print(f"[OK] {len(prs.slides)} slides -> {OUTPUT}")


if __name__ == "__main__":
    build_deck()
